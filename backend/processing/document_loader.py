from dataclasses import dataclass
from pathlib import Path

import pymupdf4llm
from docx import Document as DocxDocument
from pptx import Presentation

from backend.processing.table_extractor import extract_tables_as_markdown

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx"}


@dataclass
class LoadedPage:
    text: str
    page: int | None = None


def load_document(file_path: str) -> list[LoadedPage]:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _load_pdf(file_path)
    if ext == ".docx":
        return _load_docx(file_path)
    if ext == ".txt":
        return _load_txt(file_path)
    if ext == ".pptx":
        return _load_pptx(file_path)
    raise ValueError(f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTENSIONS)}")


def _load_pdf(file_path: str) -> list[LoadedPage]:
    pages_md = pymupdf4llm.to_markdown(file_path, page_chunks=True)
    pages = [
        LoadedPage(text=page["text"], page=page.get("metadata", {}).get("page", i + 1))
        for i, page in enumerate(pages_md)
        if page.get("text", "").strip()
    ]

    for table in extract_tables_as_markdown(file_path):
        pages.append(LoadedPage(text=table["markdown"], page=table["page"]))

    return pages


def _load_docx(file_path: str) -> list[LoadedPage]:
    doc = DocxDocument(file_path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        rows = ["| " + " | ".join(cell.text.strip() for cell in row.cells) + " |" for row in table.rows]
        if rows:
            text += "\n\n" + "\n".join(rows)
    return [LoadedPage(text=text, page=None)] if text.strip() else []


def _load_txt(file_path: str) -> list[LoadedPage]:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return [LoadedPage(text=text, page=None)] if text.strip() else []


def _load_pptx(file_path: str) -> list[LoadedPage]:
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            pages.append(LoadedPage(text="\n".join(texts), page=i))
    return pages
